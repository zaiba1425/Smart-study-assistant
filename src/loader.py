# loader.py
import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
# CORRECTED IMPORT: Document is now in langchain_core
from langchain_core.documents import Document as LangChainDocument

def load_text_from_file(filepath):
    """Loads raw text from PDF, DOCX, or PPTX files."""
    ext = os.path.splitext(filepath)[1].lower()
    text = ""

    # Handling PDF files
    if ext == ".pdf":
        try:
            reader = PdfReader(filepath)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        except Exception as e:
            # Handle potential errors during PDF reading
            text = f"Error reading PDF: {e}" 

    # Handling DOCX files
    elif ext == ".docx":
        try:
            doc = Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            text = f"Error reading DOCX: {e}"

    # Handling PPTX files
    elif ext == ".pptx":
        try:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            text = f"Error reading PPTX: {e}"
    
    # Return None if text is empty after stripping whitespace
    if not text.strip():
        return None
        
    return text

def categorize_text(text):
    """Categorizes text content based on keywords."""
    text_lower = text.lower()
    
    # Category 1: Questions/Assessment
    if any(word in text_lower for word in ["question", "quiz", "mcq", "exam", "test", "assessment"]):
        return "questions"
        
    # Category 2: Practical/Hands-on
    elif any(word in text_lower for word in ["experiment", "lab", "practical", "procedure", "data analysis"]):
        return "labs"
        
    # Category 3: General Theory/Structure
    elif any(word in text_lower for word in ["chapter", "topic", "lesson", "definition", "introduction", "summary"]):
        return "textbook"
        
    # Final Fallback
    else:
        return "notes"
        
def process_file_to_document(filepath):
    """Loads file, categorizes it, and returns a single LangChain Document."""
    raw_text = load_text_from_file(filepath)
    
    if raw_text is None:
        return None
        
    category = categorize_text(raw_text)
    
    # Create a LangChain Document object with metadata
    doc = LangChainDocument(
        page_content=raw_text,
        metadata={
            "source": os.path.basename(filepath),
            "category": category
        }
    )
    return doc